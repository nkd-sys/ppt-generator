import streamlit as st
import requests
import json
from pptx import Presentation
from pptx.util import Inches
import os

# OpenRouter API Key
API_KEY = ""  # Replace with your actual OpenRouter API key
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

# Function to generate PPT
def generate_ppt(topic, slides_content):
    prs = Presentation()
    for slide_data in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content layout
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = slide_data["title"]
        content.text = slide_data["content"]
    
    ppt_filename = f"{topic.replace(' ', '_')}.pptx"
    prs.save(ppt_filename)
    return ppt_filename

# Function to get content from OpenRouter
def get_ppt_content(topic):
    prompt = f"Create a PowerPoint outline for the topic: {topic}. Provide 8 slides with titles and key points and some images or video's."
    
    response = requests.post(
        OPENROUTER_URL,
        headers={
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        },
        data=json.dumps({
            "model": "openai/gpt-4o-mini",
            "messages": [{"role": "user", "content": prompt}]
        })
    )

    if response.status_code == 200:
        result = response.json()
        text = result["choices"][0]["message"]["content"]
        
        # Parse the AI response into slide content
        slides = []
        for slide_text in text.split("\n\n"):  # Assuming slides are separated by double line breaks
            parts = slide_text.split("\n", 1)
            if len(parts) == 2:
                slides.append({"title": parts[0].strip(), "content": parts[1].strip()})
        
        return slides if slides else None
    else:
        st.error("Error generating slide content. Please try again.")
        return None

# Streamlit UI
st.title("PowerPoint Generator 🎤📊")
st.write("Enter a topic and generate a PowerPoint presentation!")

topic = st.text_input("Enter PowerPoint Topic:")

if st.button("Generate PowerPoint"):
    if topic:
        with st.spinner("Generating"):
            slides_content = get_ppt_content(topic)
            
            if slides_content:
                ppt_file = generate_ppt(topic, slides_content)
                
                with open(ppt_file, "rb") as file:
                    btn = st.download_button(
                        label="Download PowerPoint 📥",
                        data=file,
                        file_name=ppt_file,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                
                os.remove(ppt_file)  # Clean up file after download
            else:
                st.error("Failed to generate slides. Try again.")
    else:
        st.warning("Please enter a topic.")

